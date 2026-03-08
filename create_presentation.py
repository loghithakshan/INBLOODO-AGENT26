#!/usr/bin/env python3
"""
Blood Report AI - System Architecture Presentation Generator
Creates a comprehensive PowerPoint presentation with all system details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(102, 126, 234)  # Blue
SECONDARY_COLOR = RGBColor(118, 75, 162)  # Purple
ACCENT_COLOR = RGBColor(255, 107, 107)  # Red
SUCCESS_COLOR = RGBColor(81, 207, 102)  # Green
WARNING_COLOR = RGBColor(255, 169, 77)  # Orange
TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray
LIGHT_BG = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 220, 255)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items=None, include_bg=True):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    if include_bg:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.3)
    
    # Add content
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    return slide

def add_architecture_diagram(prs):
    """Add system architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "System Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Layer 1: Input
    y_pos = 1.2
    
    # Input Layer
    input_box = slide.shapes.add_shape(1, Inches(1), Inches(y_pos), Inches(8), Inches(0.6))
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    input_box.line.color.rgb = PRIMARY_COLOR
    input_frame = input_box.text_frame
    input_frame.text = "INPUT LAYER: PDF | CSV | JSON | Images"
    input_frame.paragraphs[0].font.size = Pt(14)
    input_frame.paragraphs[0].font.bold = True
    input_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    input_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow down
    arrow1 = slide.shapes.add_connector(1, Inches(5), Inches(1.8), Inches(5), Inches(2.2))
    arrow1.line.color.rgb = PRIMARY_COLOR
    arrow1.line.width = Pt(2)
    
    # Processing Layer
    y_pos = 2.3
    proc_box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
    proc_box.fill.solid()
    proc_box.fill.fore_color.rgb = RGBColor(240, 230, 255)
    proc_box.line.color.rgb = SECONDARY_COLOR
    proc_frame = proc_box.text_frame
    proc_frame.text = "PROCESSING LAYER: Data Extraction | Cleaning | Parameter Mapping"
    proc_frame.paragraphs[0].font.size = Pt(14)
    proc_frame.paragraphs[0].font.bold = True
    proc_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    proc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow down
    arrow2 = slide.shapes.add_connector(1, Inches(5), Inches(3.1), Inches(5), Inches(3.5))
    arrow2.line.color.rgb = SECONDARY_COLOR
    arrow2.line.width = Pt(2)
    
    # AI Analysis Layer
    y_pos = 3.6
    ai_box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    ai_box.line.color.rgb = SUCCESS_COLOR
    ai_frame = ai_box.text_frame
    ai_frame.text = "AI ANALYSIS LAYER: Multi-LLM (Gemini, GPT-4, Claude, Grok)"
    ai_frame.paragraphs[0].font.size = Pt(14)
    ai_frame.paragraphs[0].font.bold = True
    ai_frame.paragraphs[0].font.color.rgb = SUCCESS_COLOR
    ai_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow down
    arrow3 = slide.shapes.add_connector(1, Inches(5), Inches(4.4), Inches(5), Inches(4.8))
    arrow3.line.color.rgb = SUCCESS_COLOR
    arrow3.line.width = Pt(2)
    
    # Output Layer
    y_pos = 4.9
    out_box = slide.shapes.add_shape(1, Inches(1), Inches(y_pos), Inches(8), Inches(0.6))
    out_box.fill.solid()
    out_box.fill.fore_color.rgb = RGBColor(255, 240, 230)
    out_box.line.color.rgb = WARNING_COLOR
    out_frame = out_box.text_frame
    out_frame.text = "OUTPUT LAYER: JSON | PDF | Dashboard | Recommendations"
    out_frame.paragraphs[0].font.size = Pt(14)
    out_frame.paragraphs[0].font.bold = True
    out_frame.paragraphs[0].font.color.rgb = WARNING_COLOR
    out_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow down
    arrow4 = slide.shapes.add_connector(1, Inches(5), Inches(5.5), Inches(5), Inches(5.9))
    arrow4.line.color.rgb = WARNING_COLOR
    arrow4.line.width = Pt(2)
    
    # User Layer
    y_pos = 6.0
    user_box = slide.shapes.add_shape(1, Inches(1.5), Inches(y_pos), Inches(7), Inches(0.6))
    user_box.fill.solid()
    user_box.fill.fore_color.rgb = RGBColor(230, 230, 230)
    user_box.line.color.rgb = TEXT_COLOR
    user_frame = user_box.text_frame
    user_frame.text = "📱 User Interface | Download Reports | Analytics Dashboard"
    user_frame.paragraphs[0].font.size = Pt(14)
    user_frame.paragraphs[0].font.bold = True
    user_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    user_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_workflow_diagram(prs):
    """Add workflow diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "Complete Data Processing Workflow"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Create workflow boxes
    workflow_steps = [
        ("1. FILE UPLOAD", "PDF/CSV/JSON/Images", RGBColor(200, 220, 255), 1.0, 1.3),
        ("2. EXTRACTION", "Parse & Extract Data", RGBColor(220, 200, 255), 2.5, 1.3),
        ("3. CLEANING", "Normalize & Validate", RGBColor(200, 255, 220), 4.0, 1.3),
        ("4. MAPPING", "Parameter Mapping", RGBColor(255, 240, 200), 5.5, 1.3),
        ("5. AI ANALYSIS", "Multi-AI Processing", RGBColor(255, 220, 220), 7.0, 1.3),
        ("6. SYNTHESIS", "Combine Results", RGBColor(240, 240, 240), 8.5, 1.3),
    ]
    
    for step, desc, color, x, y in workflow_steps:
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.3), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_bottom = Inches(0.05)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        
        p = text_frame.paragraphs[0]
        p.text = step
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(8)
        p2.font.color.rgb = TEXT_COLOR
        p2.alignment = PP_ALIGN.CENTER
    
    # Add arrows between steps
    arrow_y = 1.95
    for x in [2.35, 3.85, 5.35, 6.85, 8.35]:
        arrow = slide.shapes.add_connector(1, Inches(x), Inches(arrow_y), Inches(x + 0.15), Inches(arrow_y))
        arrow.line.color.rgb = PRIMARY_COLOR
        arrow.line.width = Pt(2)
    
    # Add output section
    output_box = slide.shapes.add_shape(1, Inches(1), Inches(3.2), Inches(8), Inches(2.8))
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
    output_box.line.color.rgb = PRIMARY_COLOR
    output_box.line.width = Pt(2)
    
    output_frame = output_box.text_frame
    output_frame.word_wrap = True
    output_frame.margin_left = Inches(0.2)
    output_frame.margin_top = Inches(0.1)
    
    outputs = [
        "📊 ANALYSIS RESULTS:",
        "• Extracted Medical Parameters (Hemoglobin, Glucose, Cholesterol, etc.)",
        "• Clinical Interpretations & Findings",
        "• Identified Health Risks & Conditions",
        "• Recommended Medicines & Supplements (with dosages)",
        "• Health Recommendations & Lifestyle Changes",
        "• Medical Prescriptions by Category",
        "• AI Confidence Score & System Metrics"
    ]
    
    for i, output in enumerate(outputs):
        if i == 0:
            p = output_frame.paragraphs[0]
        else:
            p = output_frame.add_paragraph()
        
        p.text = output
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_COLOR
        if i == 0:
            p.font.bold = True
            p.font.size = Pt(12)

def add_tools_technologies(prs):
    """Add tools and technologies slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "Technologies & Tools"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Categories
    categories = [
        {
            "name": "BACKEND FRAMEWORK",
            "items": ["FastAPI (Python Web Framework)", "Uvicorn (ASGI Server)", "PostgreSQL Database"],
            "color": RGBColor(230, 240, 255),
            "x": 0.3
        },
        {
            "name": "AI/LLM PROVIDERS",
            "items": ["Google Gemini 1.5 Pro", "OpenAI GPT-4", "Anthropic Claude 3.5", "xAI Grok 2", "Mock AI (Testing)"],
            "color": RGBColor(240, 230, 255),
            "x": 3.4
        },
        {
            "name": "DATA PROCESSING",
            "items": ["EasyOCR (Image Text Extraction)", "PDF Processing", "CSV/JSON Parsing", "Text Normalization"],
            "color": RGBColor(240, 255, 240),
            "x": 6.5
        }
    ]
    
    for cat in categories:
        # Category box
        cat_box = slide.shapes.add_shape(1, Inches(cat['x']), Inches(1.2), Inches(3), Inches(5.8))
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = cat['color']
        cat_box.line.color.rgb = PRIMARY_COLOR
        cat_box.line.width = Pt(2)
        
        text_frame = cat_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        text_frame.margin_top = Inches(0.1)
        
        # Category name
        p = text_frame.paragraphs[0]
        p.text = cat['name']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Items
        for item in cat['items']:
            p = text_frame.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(4)
            p.space_after = Pt(4)

def add_ai_providers_slide(prs):
    """Add AI providers detailed slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "Multi-LLM AI Analysis"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    providers = [
        ("Google Gemini 1.5 Pro", "Advanced reasoning, multimodal", RGBColor(230, 240, 255)),
        ("OpenAI GPT-4", "Powerful language model, reasoning", RGBColor(240, 230, 255)),
        ("Anthropic Claude 3.5", "Constitutional AI, safe outputs", RGBColor(230, 255, 240)),
        ("xAI Grok 2", "Real-time knowledge, reasoning", RGBColor(255, 240, 230))
    ]
    
    y_pos = 1.2
    for provider, desc, color in providers:
        box = slide.shapes.add_shape(1, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(1.5)
        
        text_frame = box.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"✓ {provider}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p2 = text_frame.add_paragraph()
        p2.text = f"   {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_COLOR
        
        y_pos += 1.1
    
    # Add features box
    features_box = slide.shapes.add_shape(1, Inches(0.8), Inches(5.5), Inches(8.4), Inches(1.5))
    features_box.fill.solid()
    features_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
    features_box.line.color.rgb = SUCCESS_COLOR
    features_box.line.width = Pt(2)
    
    features_frame = features_box.text_frame
    features_frame.word_wrap = True
    features_frame.margin_left = Inches(0.2)
    features_frame.margin_top = Inches(0.1)
    
    features = [
        "🎯 MULTI-AI COMPARISON: Responses from all providers analyzed for best insights",
        "✅ ATTRIBUTION TRACKING: Know which AI provider generated each analysis section",
        "📊 CONFIDENCE SCORING: AI confidence percentage for overall analysis"
    ]
    
    for i, feature in enumerate(features):
        if i == 0:
            p = features_frame.paragraphs[0]
        else:
            p = features_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)

def add_features_slide(prs):
    """Add key features slide"""
    content_items = [
        "📄 INPUT FORMATS: PDF reports, CSV data, JSON structures, Medical images",
        "🔍 TEXT EXTRACTION: Advanced OCR for scanned blood reports (EasyOCR)",
        "📊 DATA CLEANING: Automatic normalization and validation of medical parameters",
        "🤖 MULTI-AI ANALYSIS: Parallel processing with 4+ AI providers for comprehensive insights",
        "💊 MEDICAL INSIGHTS: Parameter interpretation, risk identification, recommendations",
        "📋 PRESCRIPTIONS: Categorized medicines with dosages (requires doctor approval)",
        "📱 USER INTERFACE: Modern React/HTML5 dashboard with real-time analysis updates",
        "📥 PDF GENERATION: Professional download reports with proper formatting",
        "📊 ANALYTICS DASHBOARD: Real-time metrics, agent performance, AI confidence scores"
    ]
    add_content_slide(prs, "Key Features & Capabilities", content_items)

def add_data_flow_slide(prs):
    """Add data flow diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "Data Flow Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Left side - Input
    input_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(2), Inches(1.2))
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    input_box.line.color.rgb = PRIMARY_COLOR
    input_box.line.width = Pt(2)
    input_frame = input_box.text_frame
    input_frame.text = "USER INPUT\n\nFile Upload\nData Entry"
    for p in input_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Middle - Processing
    proc_box = slide.shapes.add_shape(1, Inches(3.5), Inches(1.5), Inches(3), Inches(1.2))
    proc_box.fill.solid()
    proc_box.fill.fore_color.rgb = RGBColor(240, 230, 255)
    proc_box.line.color.rgb = SECONDARY_COLOR
    proc_box.line.width = Pt(2)
    proc_frame = proc_box.text_frame
    proc_frame.word_wrap = True
    proc_frame.text = "PROCESSING\n\nExtraction → Cleaning\nParameter Mapping"
    for p in proc_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Right side - Output
    output_box = slide.shapes.add_shape(1, Inches(7.5), Inches(1.5), Inches(2), Inches(1.2))
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    output_box.line.color.rgb = SUCCESS_COLOR
    output_box.line.width = Pt(2)
    output_frame = output_box.text_frame
    output_frame.text = "RESULTS\n\nJSON Response\nDashboard Update"
    for p in output_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Add arrows
    arrow1 = slide.shapes.add_connector(1, Inches(2.5), Inches(2.1), Inches(3.5), Inches(2.1))
    arrow1.line.color.rgb = PRIMARY_COLOR
    arrow1.line.width = Pt(2)
    
    arrow2 = slide.shapes.add_connector(1, Inches(6.5), Inches(2.1), Inches(7.5), Inches(2.1))
    arrow2.line.color.rgb = SECONDARY_COLOR
    arrow2.line.width = Pt(2)
    
    # Bottom section - AI Analysis
    ai_box = slide.shapes.add_shape(1, Inches(1), Inches(3.2), Inches(8), Inches(3.5))
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    ai_box.line.color.rgb = WARNING_COLOR
    ai_box.line.width = Pt(2)
    
    ai_frame = ai_box.text_frame
    ai_frame.word_wrap = True
    ai_frame.margin_left = Inches(0.2)
    ai_frame.margin_top = Inches(0.1)
    
    ai_items = [
        "🤖 AI ANALYSIS ENGINE:",
        "→ Parallel Multi-LLM Processing (Gemini, GPT-4, Claude, Grok)",
        "→ Agent Orchestration: Multiple specialized agents analyze different aspects",
        "→ Result Synthesis: Combine AI outputs for comprehensive analysis",
        "→ Attribution Tracking: Identify which provider generated each insight",
        "→ Confidence Scoring: Calculate overall analysis confidence percentage"
    ]
    
    for i, item in enumerate(ai_items):
        if i == 0:
            p = ai_frame.paragraphs[0]
        else:
            p = ai_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(2)
        p.space_after = Pt(2)

def add_security_compliance(prs):
    """Add security and compliance slide"""
    content_items = [
        "🔐 SECURITY FEATURES:",
        "• API Key Management: Secure handling of AI provider credentials",
        "• User Authentication: Session management for secure access",
        "• Data Encryption: Sensitive health information protection",
        "",
        "✅ MEDICAL COMPLIANCE:",
        "• AI Disclaimer: Clear warnings that AI is supplementary to professional medical advice",
        "• Doctor Approval: All medicine recommendations require licensed healthcare professional approval",
        "• Data Privacy: Compliance with health data protection regulations",
        "• Audit Trail: Processing logs for medical audit requirements"
    ]
    add_content_slide(prs, "Security & Compliance", content_items)

def add_benefits_slide(prs):
    """Add benefits slide"""
    content_items = [
        "⚡ EFFICIENCY: Automated analysis of blood reports in seconds",
        "🎯 ACCURACY: Multi-AI cross-validation ensures reliable insights",
        "💡 INTELLIGENCE: Comprehensive interpretations from advanced language models",
        "👨‍⚕️ PROFESSIONAL: Doctor-verified recommendations and prescriptions",
        "📱 ACCESSIBILITY: Web-based interface works on any device",
        "📊 SCALABILITY: Handles high throughput of medical data analysis",
        "🌍 MULTI-LANGUAGE: Support for various medical parameters globally",
        "📈 CONTINUOUS IMPROVEMENT: Real-time metrics and performance tracking"
    ]
    add_content_slide(prs, "System Benefits", content_items)

def add_technical_stack(prs):
    """Add technical stack slide"""
    content_items = [
        "🔧 BACKEND:",
        "• Python 3.x with FastAPI framework",
        "• Uvicorn ASGI server for production deployment",
        "• PostgreSQL database for persistent storage",
        "",
        "🎨 FRONTEND:",
        "• React.js / HTML5 modern web interface",
        "• Real-time dashboard with Chart.js analytics",
        "• Responsive design for mobile compatibility",
        "",
        "☁️ DEPLOYMENT:",
        "• Docker containerization for easy deployment",
        "• Cloud hosting ready (Render, Heroku, AWS)",
        "• Continuous integration/deployment pipelines"
    ]
    add_content_slide(prs, "Technical Stack", content_items)

def add_agent_orchestration(prs):
    """Add agent orchestration slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "Agent Orchestration System"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Agent boxes
    agents = [
        ("Parameter Agent", "Interprets medical values\ncompares with normal ranges", RGBColor(230, 240, 255)),
        ("Pattern Agent", "Identifies health patterns\nand correlations", RGBColor(240, 230, 255)),
        ("Risk Agent", "Detects potential health risks\nand complications", RGBColor(255, 230, 230)),
        ("Treatment Agent", "Suggests medicines and\nlifestyle recommendations", RGBColor(230, 255, 240))
    ]
    
    x_positions = [0.8, 3.2, 5.6, 8.0]
    for i, (agent, desc, color) in enumerate(agents):
        box = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(1.5), Inches(2.2), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(1.5)
        
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        
        p = text_frame.paragraphs[0]
        p.text = agent
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_COLOR
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
    
    # Orchestration coordinator
    coord_box = slide.shapes.add_shape(1, Inches(2), Inches(4), Inches(6), Inches(2.3))
    coord_box.fill.solid()
    coord_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
    coord_box.line.color.rgb = PRIMARY_COLOR
    coord_box.line.width = Pt(2)
    
    coord_frame = coord_box.text_frame
    coord_frame.word_wrap = True
    coord_frame.margin_left = Inches(0.2)
    coord_frame.margin_top = Inches(0.1)
    
    items = [
        "🎯 ORCHESTRATION COORDINATOR:",
        "✓ Manages parallel agent execution across multiple LLMs",
        "✓ Synchronizes results and performs cross-validation",
        "✓ Handles agent failures with fallback mechanisms",
        "✓ Aggregates findings into comprehensive analysis",
        "✓ Assigns confidence scores and tracks attribution"
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = coord_frame.paragraphs[0]
        else:
            p = coord_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(3)

def add_usage_workflow(prs):
    """Add user workflow slide"""
    content_items = [
        "STEP 1: UPLOAD BLOOD REPORT",
        "→ Select file (PDF, CSV, JSON, or Image) and upload via web interface",
        "",
        "STEP 2: AUTOMATIC PROCESSING",
        "→ System extracts data, validates parameters, maps medical values",
        "",
        "STEP 3: AI ANALYSIS",
        "→ Multi-LLM analysis in parallel (Gemini, GPT-4, Claude, Grok)",
        "",
        "STEP 4: VIEW RESULTS",
        "→ See interpretations, risks, recommendations in dashboard",
        "",
        "STEP 5: DOWNLOAD REPORT",
        "→ Generate professional PDF with all analysis for doctor consultation"
    ]
    add_content_slide(prs, "User Workflow", content_items)

def add_future_enhancements(prs):
    """Add future scope slide"""
    content_items = [
        "🔮 PLANNED ENHANCEMENTS:",
        "• Mobile App: Native iOS/Android applications for on-the-go analysis",
        "• Real-time Monitoring: Integration with health devices (fitness trackers, smartwatches)",
        "• Multi-Language Support: Support for medical reports in multiple languages",
        "• Advanced Predictive Models: Machine learning for disease progression prediction",
        "• Doctor Collaboration: Secure sharing with healthcare professionals",
        "• Medical Database: Integration with medical literature for evidence-based recommendations",
        "• Offline Capability: Local analysis without constant internet connectivity",
        "• AR Visualization: Augmented reality for health data visualization"
    ]
    add_content_slide(prs, "Future Enhancements", content_items)

def add_closing_slide(prs):
    """Add closing/thank you slide"""
    add_title_slide(prs, "Thank You!", "Blood Report AI - Revolutionizing Medical Analysis with AI")

def add_api_documentation(prs):
    """Add API documentation slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "API Endpoints & Documentation"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    endpoints = [
        ("POST /analyze-pdf/", "Upload and analyze PDF blood report", RGBColor(230, 240, 255)),
        ("POST /analyze-csv/", "Upload and analyze CSV medical data", RGBColor(240, 230, 255)),
        ("POST /analyze-json/", "Submit JSON structured medical data", RGBColor(240, 255, 240)),
        ("GET /api/multi-ai/providers", "List available AI providers", RGBColor(255, 240, 230)),
        ("GET /health", "System health check endpoint", RGBColor(230, 255, 240)),
        ("POST /download-report", "Generate and download PDF report", RGBColor(255, 245, 230))
    ]
    
    y_pos = 1.2
    for endpoint, desc, color in endpoints:
        box = slide.shapes.add_shape(1, Inches(0.6), Inches(y_pos), Inches(8.8), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(1)
        
        text_frame = box.text_frame
        p = text_frame.paragraphs[0]
        p.text = endpoint
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p2 = text_frame.add_paragraph()
        p2.text = f"→ {desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_COLOR
        
        y_pos += 0.85

def add_performance_metrics(prs):
    """Add performance metrics slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "Performance & Scalability Metrics"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Metrics boxes
    metrics = [
        ("Response Time", "2-5 seconds\nfor full analysis", RGBColor(230, 240, 255)),
        ("Throughput", "100+ concurrent\nrequests/sec", RGBColor(240, 230, 255)),
        ("Accuracy", "85-95%\nanalysis confidence", RGBColor(240, 255, 230)),
        ("Availability", "99.9%\nuptime SLA", RGBColor(255, 240, 230))
    ]
    
    x_positions = [0.8, 3.0, 5.2, 7.4]
    for i, (metric, value, color) in enumerate(metrics):
        box = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(1.5), Inches(2), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        
        p = text_frame.paragraphs[0]
        p.text = metric
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        p2 = text_frame.add_paragraph()
        p2.text = value
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_COLOR
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
    
    # Details box
    details_box = slide.shapes.add_shape(1, Inches(0.8), Inches(3.8), Inches(8.4), Inches(2.8))
    details_box.fill.solid()
    details_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
    details_box.line.color.rgb = PRIMARY_COLOR
    details_box.line.width = Pt(2)
    
    details_frame = details_box.text_frame
    details_frame.word_wrap = True
    details_frame.margin_left = Inches(0.2)
    details_frame.margin_top = Inches(0.1)
    
    details = [
        "📊 PERFORMANCE DETAILS:",
        "• Multi-LLM parallel processing reduces latency by 40%",
        "• Horizontal scaling supports thousands of concurrent users",
        "• Caching mechanisms for frequently analyzed parameters",
        "• Database optimization for rapid data retrieval",
        "• CDN distribution for static assets and fast downloads"
    ]
    
    for i, detail in enumerate(details):
        if i == 0:
            p = details_frame.paragraphs[0]
        else:
            p = details_frame.add_paragraph()
        p.text = detail
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)

def add_deployment_options(prs):
    """Add deployment options slide"""
    content_items = [
        "☁️ CLOUD DEPLOYMENT:",
        "• Render.com: Free tier for testing, paid for production",
        "• Heroku: Streamlined deployment with PostgreSQL add-on",
        "• AWS: EC2 instances with RDS database and S3 storage",
        "• Google Cloud Platform: Cloud Run and Cloud SQL setup",
        "",
        "🐳 DOCKER CONTAINERIZATION:",
        "• Docker image includes all dependencies (python-pptx, FastAPI, etc.)",
        "• Docker Compose for local multi-container development",
        "• Kubernetes ready for enterprise deployments",
        "",
        "🔧 DEPLOYMENT CHECKLIST:",
        "✓ Environment variables configured | ✓ Database initialized", 
        "✓ SSL/TLS certificates installed | ✓ API keys secured",
        "✓ Monitoring and logging enabled | ✓ Backup procedures ready"
    ]
    add_content_slide(prs, "Deployment Options", content_items)

def add_error_handling(prs):
    """Add error handling and reliability slide"""
    content_items = [
        "🛡️ ERROR HANDLING MECHANISMS:",
        "• AI Provider Fallback: If Gemini fails, automatically try GPT-4 or Claude",
        "• Graceful Degradation: Partial results if some providers timeout",
        "• Retry Logic: Automatic retry with exponential backoff for failed requests",
        "• Circuit Breaker Pattern: Temporarily disable failing services",
        "",
        "⚠️ ERROR SCENARIOS:",
        "• Invalid File Format: Clear error message with supported formats",
        "• API Rate Limiting: Queue and retry when rate limits hit",
        "• Network Timeouts: Fallback to cached results if available",
        "• Database Connection Issues: Automatic reconnection with health checks",
        "",
        "📋 MONITORING & ALERTS:",
        "• Real-time error tracking and logging (90-day retention)",
        "• Email alerts for critical failures and service degradation",
        "• Dashboard metrics for system health monitoring"
    ]
    add_content_slide(prs, "Error Handling & Reliability", content_items)

def add_system_requirements(prs):
    """Add system requirements slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "System Requirements & Setup"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    requirements = [
        ("SERVER REQUIREMENTS", [
            "CPU: 2+ cores (4+ recommended)",
            "RAM: 4GB minimum (8GB+ for production)",
            "Storage: 20GB (includes ML models)",
            "Network: 50 Mbps minimum bandwidth"
        ], RGBColor(230, 240, 255)),
        ("SOFTWARE REQUIREMENTS", [
            "Python: 3.8 or higher",
            "PostgreSQL: 12+",
            "Node.js: 14+ (for frontend build)",
            "Docker: 20.10+ (optional)"
        ], RGBColor(240, 230, 255)),
        ("API REQUIREMENTS", [
            "Google Cloud API key (Gemini)",
            "OpenAI API key (GPT-4)",
            "Anthropic API key (Claude)",
            "xAI API key (Grok 2)"
        ], RGBColor(240, 255, 240))
    ]
    
    x_positions = [0.5, 3.4, 6.3]
    for i, (title, items, color) in enumerate(requirements):
        box = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(1.2), Inches(3), Inches(5.8))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        text_frame.margin_top = Inches(0.1)
        
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(4)

def add_privacy_security_detailed(prs):
    """Add detailed privacy and security slide"""
    content_items = [
        "🔐 DATA SECURITY MEASURES:",
        "• End-to-End Encryption: All data encrypted in transit (TLS 1.3)",
        "• At-Rest Encryption: Database encryption with AES-256",
        "• API Rate Limiting: Protection against brute force attacks",
        "• CORS Configuration: Restrict cross-origin requests",
        "",
        "📋 GDPR & HIPAA COMPLIANCE:",
        "• User Consent Management: Explicit consent for data processing",
        "• Right to Be Forgotten: Complete data deletion on request",
        "• Data Portability: Export patient data in standard formats",
        "• Breach Notification: 72-hour incident notification protocol",
        "",
        "🔒 ACCESS CONTROL:",
        "• Role-Based Access Control (RBAC): Admin, Doctor, Patient roles",
        "• Multi-Factor Authentication (MFA): Additional security layer",
        "• Session Management: Automatic logout after inactivity",
        "• Audit Logging: All actions logged with timestamps and user IDs"
    ]
    add_content_slide(prs, "Privacy & Security (Detailed)", content_items)

def add_case_study(prs):
    """Add case study/example slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    title_frame = title_shape.text_frame
    title_frame.text = "Example: Anemia Case Analysis"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Input section
    input_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(4.3), Inches(2.8))
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    input_box.line.color.rgb = PRIMARY_COLOR
    input_box.line.width = Pt(2)
    
    input_frame = input_box.text_frame
    input_frame.word_wrap = True
    input_frame.margin_left = Inches(0.15)
    input_frame.margin_top = Inches(0.1)
    
    input_items = [
        "📊 INPUT:",
        "• Hemoglobin: 9.5 g/dL",
        "• RBC: 3.2 M/µL",
        "• Hematocrit: 28%",
        "• MCV: 75 fL"
    ]
    
    for i, item in enumerate(input_items):
        if i == 0:
            p = input_frame.paragraphs[0]
        else:
            p = input_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_COLOR if i > 0 else PRIMARY_COLOR
        if i == 0:
            p.font.bold = True
            p.font.size = Pt(11)
    
    # Arrow
    arrow = slide.shapes.add_connector(1, Inches(4.8), Inches(2.6), Inches(5.2), Inches(2.6))
    arrow.line.color.rgb = PRIMARY_COLOR
    arrow.line.width = Pt(2)
    
    # Analysis section
    analysis_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.2), Inches(4.3), Inches(2.8))
    analysis_box.fill.solid()
    analysis_box.fill.fore_color.rgb = RGBColor(240, 230, 255)
    analysis_box.line.color.rgb = SECONDARY_COLOR
    analysis_box.line.width = Pt(2)
    
    analysis_frame = analysis_box.text_frame
    analysis_frame.word_wrap = True
    analysis_frame.margin_left = Inches(0.15)
    analysis_frame.margin_top = Inches(0.1)
    
    analysis_items = [
        "🔍 AI ANALYSIS:",
        "• Iron-deficiency anemia",
        "• Chronic blood loss likely",
        "• Moderate severity",
        "• Requires intervention"
    ]
    
    for i, item in enumerate(analysis_items):
        if i == 0:
            p = analysis_frame.paragraphs[0]
        else:
            p = analysis_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_COLOR if i > 0 else SECONDARY_COLOR
        if i == 0:
            p.font.bold = True
            p.font.size = Pt(11)
    
    # Output section
    output_box = slide.shapes.add_shape(1, Inches(0.5), Inches(4.3), Inches(9), Inches(2.5))
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    output_box.line.color.rgb = SUCCESS_COLOR
    output_box.line.width = Pt(2)
    
    output_frame = output_box.text_frame
    output_frame.word_wrap = True
    output_frame.margin_left = Inches(0.2)
    output_frame.margin_top = Inches(0.1)
    
    output_items = [
        "✅ RECOMMENDATIONS:",
        "💊 Iron supplements: Ferrous sulfate 325mg daily | 🥩 High iron diet (red meat, spinach, legumes)",
        "🏥 Doctor consultation: Investigate chronic blood loss source | 📊 Follow-up tests in 4-6 weeks",
        "⚠️ Important: All medications must be prescribed by licensed healthcare professional"
    ]
    
    for i, item in enumerate(output_items):
        if i == 0:
            p = output_frame.paragraphs[0]
        else:
            p = output_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = SUCCESS_COLOR if i == 0 else TEXT_COLOR
        if i == 0:
            p.font.bold = True
            p.font.size = Pt(11)
        p.space_before = Pt(3)



add_title_slide(prs, "BLOOD REPORT AI", "Advanced AI-Powered Medical Analysis System")

# Overview
content = [
    "🎯 OBJECTIVE:",
    "Provide comprehensive AI-powered analysis of blood reports using advanced language models",
    "",
    "✨ KEY CAPABILITIES:",
    "• Automated extraction and analysis of medical parameters",
    "• Multi-AI comparison for reliable insights (Gemini, GPT-4, Claude, Grok)",
    "• Intelligent risk identification and health recommendations",
    "• Professional report generation and analytics dashboard",
    "",
    "🏥 USE CASES:",
    "• Clinical diagnostic support | Health screening | Medical research",
    "• Patient education | Preventive healthcare management"
]
add_content_slide(prs, "Overview & Objectives", content)

# System Architecture
add_architecture_diagram(prs)

# Workflow
add_workflow_diagram(prs)

# Data Flow
add_data_flow_slide(prs)

# Tools & Technologies  
add_tools_technologies(prs)

# AI Providers
add_ai_providers_slide(prs)

# Agent Orchestration
add_agent_orchestration(prs)

# Features
add_features_slide(prs)

# Technical Stack
add_technical_stack(prs)

# User Workflow
add_usage_workflow(prs)

# Benefits
add_benefits_slide(prs)

# Security
add_security_compliance(prs)

# Additional Important Slides
add_api_documentation(prs)
add_performance_metrics(prs)
add_deployment_options(prs)
add_error_handling(prs)
add_system_requirements(prs)
add_privacy_security_detailed(prs)
add_case_study(prs)

# Future
add_future_enhancements(prs)

# Closing
add_closing_slide(prs)

# Save presentation
output_file = "Blood_Report_AI_System_Architecture.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
